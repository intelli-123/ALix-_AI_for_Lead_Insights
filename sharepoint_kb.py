#sharepoint_kb.py
import os
import requests
import sqlite3
from datetime import datetime
from msal import ConfidentialClientApplication
from dotenv import load_dotenv
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

load_dotenv()

# === Config ===
CLIENT_ID = os.getenv("SP_CLIENT_ID")
CLIENT_SECRET = os.getenv("SP_CLIENT_SECRET")
TENANT_ID = os.getenv("SP_TENANT_ID")
SITE_DOMAIN = os.getenv("SP_SITE_DOMAIN")
# New line
raw_paths = os.getenv("SP_SITE_PATHS", "")
SITE_PATHS = [p.strip().strip("'\"") for p in raw_paths.split(',') if p.strip()]

AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
SCOPE = ["https://graph.microsoft.com/.default"]
GRAPH_API = "https://graph.microsoft.com/v1.0"
DB_FILE = "kb_cache.db" # Database file configuration

# === Ensure required folders exist ===
os.makedirs("tmp/sharepoint_docs", exist_ok=True)
os.makedirs("kb", exist_ok=True)


# === MODIFICATION START: Database functions moved here ===
def init_db():
    """Initializes the database and creates the table if it doesn't exist."""
    with sqlite3.connect(DB_FILE) as conn:
        cursor = conn.cursor()
        # Create table to store the entire KB content as a single entry
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS knowledge_base (
                id INTEGER PRIMARY KEY,
                content TEXT NOT NULL,
                last_updated TIMESTAMP NOT NULL
            )
        """)
        conn.commit()

def update_kb_in_db(content: str):
    """
    Updates the knowledge base content in the database.
    This function will replace the old KB with the new content.
    """
    with sqlite3.connect(DB_FILE) as conn:
        cursor = conn.cursor()
        # Remove the old entry
        cursor.execute("DELETE FROM knowledge_base")
        # Insert the new entry with the current timestamp
        cursor.execute(
            "INSERT INTO knowledge_base (content, last_updated) VALUES (?, ?)",
            (content, datetime.now())
        )
        conn.commit()
        print("✅ Database cache has been updated.")

def get_kb_from_db() -> str | None:
    """
    Retrieves the knowledge base content from the database.
    Returns the content as a string, or None if the cache is empty.
    """
    try:
        with sqlite3.connect(DB_FILE) as conn:
            cursor = conn.cursor()
            # Get the most recent entry
            cursor.execute("SELECT content FROM knowledge_base ORDER BY last_updated DESC LIMIT 1")
            row = cursor.fetchone()
            if row:
                print("✅ KB content retrieved from database cache.")
                return row[0]
            else:
                print("🟡 Database cache is empty.")
                return None
    except sqlite3.OperationalError:
        # This can happen if the table doesn't exist yet
        print("🟡 Database table not found. It will be created.")
        return None
# === MODIFICATION END ===


# === File extractors ===
def extract_text_from_pdf(filepath):
    try:
        text = ""
        with fitz.open(filepath) as doc:
            for page in doc:
                text += page.get_text()
        return text
    except Exception as e:
        print(f"❌ Error extracting {filepath}: {e}")
        return ""

def extract_text_from_docx(filepath):
    try:
        doc = Document(filepath)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        print(f"❌ Error extracting {filepath}: {e}")
        return ""

def extract_text_from_pptx(filepath):
    try:
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"❌ Error extracting {filepath}: {e}")
        return ""

def extract_text_from_txt(filepath):
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        print(f"❌ Error extracting {filepath}: {e}")
        return ""

# === Recursive traversal of folders (up to 2 levels) ===
def traverse_folder(site_id, item_id, headers, level=0, max_depth=2):
    if level > max_depth:
        return ""

    combined = ""
    children_url = f"{GRAPH_API}/sites/{site_id}/drive/items/{item_id}/children"
    res = requests.get(children_url, headers=headers)

    if res.status_code != 200:
        print(f"⚠️ Failed to list children: {res.text}")
        return ""

    for item in res.json().get("value", []):
        if "folder" in item:
            combined += traverse_folder(site_id, item["id"], headers, level + 1, max_depth)
        elif "file" in item:
            name = item["name"]
            file_id = item["id"]
            ext = os.path.splitext(name)[1].lower()
            if ext in [".pdf", ".docx", ".pptx", ".txt"]:
                print(f"📁📁📁 Downloading: {name}")
                dl_url = f"{GRAPH_API}/sites/{site_id}/drive/items/{file_id}/content"
                res_file = requests.get(dl_url, headers=headers)

                if res_file.status_code == 200:
                    local_path = os.path.join("tmp/sharepoint_docs", f"{site_id}_{name}")
                    with open(local_path, "wb") as f:
                        f.write(res_file.content)

                    # Extract content
                    if ext == ".pdf":
                        extracted = extract_text_from_pdf(local_path)
                    elif ext == ".docx":
                        extracted = extract_text_from_docx(local_path)
                    elif ext == ".pptx":
                        extracted = extract_text_from_pptx(local_path)
                    elif ext == ".txt":
                        extracted = extract_text_from_txt(local_path)
                    else:
                        extracted = ""

                    if extracted.strip():
                        print(f"✅ Extracted from: {name}")
                        combined += f"\n\n# Document: {name}\n{extracted.strip()}"
                    else:
                        print(f"⚠️ No content extracted from: {name}")
                else:
                    print(f"❌ Failed to download {name} ({res_file.status_code})")
            else:
                print(f"⏭️ Skipping unsupported file: {name}")

    return combined

# === Main function ===
def get_sharepoint_kb():
    print("=== Starting SharePoint KB Extraction ===")
    print("🔑 Authenticating to Microsoft Graph...")

    app = ConfidentialClientApplication(
        client_id=CLIENT_ID,
        authority=AUTHORITY,
        client_credential=CLIENT_SECRET
    )
    token_response = app.acquire_token_for_client(scopes=SCOPE)

    if "access_token" not in token_response:
        print("❌ SP Authentication failed.")
        return ""

    token = token_response["access_token"]
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json"
    }

    print(f"✅ SITE_PATHS: {SITE_PATHS}")
    all_text = ""

    for path in SITE_PATHS:
        path = path.strip()
        print(f"\n🔍 Processing site path: {path}")
        site_url = f"{GRAPH_API}/sites/{SITE_DOMAIN}:{path}"
        res = requests.get(site_url, headers=headers)

        if res.status_code != 200:
            print(f"❌ Site not found for path: {path}")
            print("Response:", res.text)
            continue

        site_id = res.json()["id"]
        print(f"✅ Found site ID: {site_id}")

        # Root folder
        root_items_url = f"{GRAPH_API}/sites/{site_id}/drive/root/children"
        root_res = requests.get(root_items_url, headers=headers)

        if root_res.status_code != 200:
            print(f"❌ Could not list root drive contents for {path}")
            continue

        for item in root_res.json().get("value", []):
            if "folder" in item or "file" in item:
                extracted = traverse_folder(site_id, item["id"], headers, level=1)
                all_text += extracted

    return all_text.strip()


# === Entry point now saves to DB instead of file ===
if __name__ == "__main__":
    # 1. Initialize the database to ensure the table exists
    print("🚀 Initializing database...")
    init_db()

    # 2. Fetch the latest content from SharePoint
    kb_content = get_sharepoint_kb()

    # 3. If content was fetched, update it in the database
    if kb_content:
        print("\n💾 Saving knowledge base content to the database cache...")
        update_kb_in_db(kb_content)
        print("✅ Knowledge base has been successfully updated in the database.")
    else:
        print("⚠️ No new content was extracted from SharePoint. Database remains unchanged.")



# # sharepoint_kb.py

# import os
# import requests
# from msal import ConfidentialClientApplication
# from dotenv import load_dotenv
# import fitz  # PyMuPDF
# from docx import Document
# from pptx import Presentation

# load_dotenv()

# # === Config ===
# CLIENT_ID = os.getenv("SP_CLIENT_ID")
# CLIENT_SECRET = os.getenv("SP_CLIENT_SECRET")
# TENANT_ID = os.getenv("SP_TENANT_ID")
# SITE_DOMAIN = os.getenv("SP_SITE_DOMAIN")
# SITE_PATHS = os.getenv("SP_SITE_PATHS", "").split(",")  # Comma-separated paths

# AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
# SCOPE = ["https://graph.microsoft.com/.default"]
# GRAPH_API = "https://graph.microsoft.com/v1.0"

# # === Ensure required folders exist ===
# os.makedirs("tmp/sharepoint_docs", exist_ok=True)
# os.makedirs("kb", exist_ok=True)

# # === File extractors ===
# def extract_text_from_pdf(filepath):
#     try:
#         text = ""
#         with fitz.open(filepath) as doc:
#             for page in doc:
#                 text += page.get_text()
#         return text
#     except Exception as e:
#         print(f"❌ Error extracting {filepath}: {e}")
#         return ""

# def extract_text_from_docx(filepath):
#     try:
#         doc = Document(filepath)
#         return "\n".join([p.text for p in doc.paragraphs])
#     except Exception as e:
#         print(f"❌ Error extracting {filepath}: {e}")
#         return ""

# def extract_text_from_pptx(filepath):
#     try:
#         prs = Presentation(filepath)
#         text = ""
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text + "\n"
#         return text
#     except Exception as e:
#         print(f"❌ Error extracting {filepath}: {e}")
#         return ""

# def extract_text_from_txt(filepath):
#     try:
#         with open(filepath, "r", encoding="utf-8") as f:
#             return f.read()
#     except Exception as e:
#         print(f"❌ Error extracting {filepath}: {e}")
#         return ""

# # === Recursive traversal of folders (up to 2 levels) ===
# def traverse_folder(site_id, item_id, headers, level=0, max_depth=2):
#     if level > max_depth:
#         return ""

#     combined = ""
#     children_url = f"{GRAPH_API}/sites/{site_id}/drive/items/{item_id}/children"
#     res = requests.get(children_url, headers=headers)

#     if res.status_code != 200:
#         print(f"⚠️ Failed to list children: {res.text}")
#         return ""

#     for item in res.json().get("value", []):
#         if "folder" in item:
#             combined += traverse_folder(site_id, item["id"], headers, level + 1, max_depth)
#         elif "file" in item:
#             name = item["name"]
#             file_id = item["id"]
#             ext = os.path.splitext(name)[1].lower()
#             if ext in [".pdf", ".docx", ".pptx", ".txt"]:
#                 print(f"📁📁📁 Downloading: {name}")
#                 dl_url = f"{GRAPH_API}/sites/{site_id}/drive/items/{file_id}/content"
#                 res_file = requests.get(dl_url, headers=headers)

#                 if res_file.status_code == 200:
#                     local_path = os.path.join("tmp/sharepoint_docs", f"{site_id}_{name}")
#                     with open(local_path, "wb") as f:
#                         f.write(res_file.content)

#                     # Extract content
#                     if ext == ".pdf":
#                         extracted = extract_text_from_pdf(local_path)
#                     elif ext == ".docx":
#                         extracted = extract_text_from_docx(local_path)
#                     elif ext == ".pptx":
#                         extracted = extract_text_from_pptx(local_path)
#                     elif ext == ".txt":
#                         extracted = extract_text_from_txt(local_path)
#                     else:
#                         extracted = ""

#                     if extracted.strip():
#                         print(f"✅ Extracted from: {name}")
#                         combined += f"\n\n# Document: {name}\n{extracted.strip()}"
#                     else:
#                         print(f"⚠️ No content extracted from: {name}")
#                 else:
#                     print(f"❌ Failed to download {name} ({res_file.status_code})")
#             else:
#                 print(f"⏭️ Skipping unsupported file: {name}")

#     return combined

# # === Main function ===
# def get_sharepoint_kb():
#     print("=== Starting SharePoint KB Extraction ===")
#     print("🔑 Authenticating to Microsoft Graph...")

#     app = ConfidentialClientApplication(
#         client_id=CLIENT_ID,
#         authority=AUTHORITY,
#         client_credential=CLIENT_SECRET
#     )
#     token_response = app.acquire_token_for_client(scopes=SCOPE)

#     if "access_token" not in token_response:
#         print("❌ SP Authentication failed.")
#         return ""

#     token = token_response["access_token"]
#     headers = {
#         "Authorization": f"Bearer {token}",
#         "Content-Type": "application/json"
#     }

#     print(f"✅ SITE_PATHS: {SITE_PATHS}")
#     all_text = ""

#     for path in SITE_PATHS:
#         path = path.strip()
#         print(f"\n🔍 Processing site path: {path}")
#         site_url = f"{GRAPH_API}/sites/{SITE_DOMAIN}:{path}"
#         res = requests.get(site_url, headers=headers)

#         if res.status_code != 200:
#             print(f"❌ Site not found for path: {path}")
#             print("Response:", res.text)
#             continue

#         site_id = res.json()["id"]
#         print(f"✅ Found site ID: {site_id}")

#         # Root folder
#         root_items_url = f"{GRAPH_API}/sites/{site_id}/drive/root/children"
#         root_res = requests.get(root_items_url, headers=headers)

#         if root_res.status_code != 200:
#             print(f"❌ Could not list root drive contents for {path}")
#             continue

#         for item in root_res.json().get("value", []):
#             if "folder" in item or "file" in item:
#                 extracted = traverse_folder(site_id, item["id"], headers, level=1)
#                 all_text += extracted

#     return all_text.strip()


# # === Entry point ===
# if __name__ == "__main__":
#     kb_content = get_sharepoint_kb()

#     if kb_content:
#         with open("kb/sp_cache.md", "w", encoding="utf-8") as f:
#             f.write(kb_content)
#         print("✅ Knowledge base saved to: kb/sp_cache.md")
#     else:
#         print("⚠️ No content extracted.")


# # sharepoint_kb.py

# import os, requests
# from msal import ConfidentialClientApplication
# from dotenv import load_dotenv

# load_dotenv()

# CLIENT_ID = os.getenv("SP_CLIENT_ID")
# CLIENT_SECRET = os.getenv("SP_CLIENT_SECRET")
# TENANT_ID = os.getenv("SP_TENANT_ID")
# SITE_DOMAIN = os.getenv("SP_SITE_DOMAIN")
# SITE_PATH = os.getenv("SP_SITE_PATH")

# AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
# SCOPE = ["https://graph.microsoft.com/.default"]
# GRAPH_API = "https://graph.microsoft.com/v1.0"

# def get_sharepoint_kb():
#     app = ConfidentialClientApplication(
#         client_id=CLIENT_ID,
#         authority=AUTHORITY,
#         client_credential=CLIENT_SECRET
#     )
#     token_response = app.acquire_token_for_client(scopes=SCOPE)

#     if 'access_token' not in token_response:
#         print("❌ SP Authentication failed.")
#         return ""

#     token = token_response['access_token']
#     headers = {
#         'Authorization': f'Bearer {token}',
#         'Content-Type': 'application/json'
#     }

#     response = requests.get(
#         f'{GRAPH_API}/sites/{SITE_DOMAIN}:{SITE_PATH}',
#         headers=headers
#     )

#     if response.status_code == 200:
#         site = response.json()
#         site_id = site['id']

#         # Get pages or documents (you can enhance this later)
#         drive_items = requests.get(f'{GRAPH_API}/sites/{site_id}/drive/root/children', headers=headers)
#         if drive_items.status_code == 200:
#             kb_text = ""
#             for item in drive_items.json().get('value', []):
#                 if 'file' in item:
#                     kb_text += f"\n[File: {item['name']}] - {item.get('webUrl')}\n"
#             return kb_text
#         else:
#             print("❌ Failed to fetch drive items.")
#             return ""
#     else:
#         print(f"❌ SP Site not found. Status: {response.status_code}")
#         print("URL used:", f'{GRAPH_API}/sites/{SITE_DOMAIN}:{SITE_PATH}')
#         print("Response body:", response.text)



# import requests
# from msal import ConfidentialClientApplication

# # Azure AD App credentials
# CLIENT_ID = 'f843b4db-1107-40d3-bd23-bacfda35cc29'
# CLIENT_SECRET = 'EWH8Q~8AHB_GEXesLBIaoC884kCx9oXZlNkwAdBj'
# TENANT_ID = '89cf11d4-079d-47a6-af93-e6ae64ceb42c'

# # Microsoft Graph API endpoint
# AUTHORITY = f"https://login.microsoftonline.com/{TENANT_ID}"
# SCOPE = ["https://graph.microsoft.com/.default"]
# GRAPH_API = "https://graph.microsoft.com/v1.0"


# # Create a confidential client app
# app = ConfidentialClientApplication(
#     client_id=CLIENT_ID,
#     authority=AUTHORITY,
#     client_credential=CLIENT_SECRET
# )

# # Acquire token
# token_response = app.acquire_token_for_client(scopes=SCOPE)

# if 'access_token' in token_response:
#     token = token_response['access_token']
#     headers = {
#         'Authorization': f'Bearer {token}',
#         'Content-Type': 'application/json'
#     }
# #https://intelliswift.sharepoint.com/sites/TechnologyCommunities/DevOps
#     # Query a specific SharePoint site by path
#     site_url = 'intelliswift.sharepoint.com'      # 🔁 Replace with your actual domain
#     site_path = '/sites/TechnologyCommunities/DevOps'              # 🔁 Replace with your actual site path

#     response = requests.get(
#         f'{GRAPH_API}/sites/{site_url}:{site_path}',
#         headers=headers
#     )

#     if response.status_code == 200:
#         site = response.json()
#         print(f"✅ Found site: {site['name']} - URL: {site['webUrl']} - ID: {site['id']}")
#     else:
#         print(f"❌ Site not found. Status: {response.status_code}")
#         print(response.text)

# else:
#     print("❌ Authentication failed.")
#     print(token_response.get("error_description"))
